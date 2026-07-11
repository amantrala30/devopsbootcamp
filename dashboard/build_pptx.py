"""
RHP Modernization – Executive Status Presentation Builder
Generates dashboard/RHP_Executive_Status.pptx from project status report data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import copy

# ── Palette (Light Theme) ─────────────────────────────────────────────────────
BG_DARK    = RGBColor(0xF0, 0xF6, 0xFF)   # very light blue page background
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)   # white card
SURFACE2   = RGBColor(0xE5, 0xEE, 0xFC)   # light blue-gray alternate rows
BORDER     = RGBColor(0xC4, 0xD2, 0xEB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT       = RGBColor(0x0F, 0x17, 0x2A)   # near-black primary text
MUTED      = RGBColor(0x4A, 0x5A, 0x82)   # medium blue-gray secondary text
ACCENT     = RGBColor(0x25, 0x63, 0xEB)   # strong blue
GREEN      = RGBColor(0x16, 0x80, 0x3D)   # dark green (readable on light)
YELLOW     = RGBColor(0xD9, 0x77, 0x06)   # amber (readable on light)
RED        = RGBColor(0xDC, 0x26, 0x26)   # red
PURPLE     = RGBColor(0x6D, 0x28, 0xD9)   # purple
GREEN_DIM  = RGBColor(0xDC, 0xFC, 0xE7)   # light green tint
YELLOW_DIM = RGBColor(0xFE, 0xF3, 0xC7)   # light yellow tint
RED_DIM    = RGBColor(0xFE, 0xE2, 0xE2)   # light red tint
BLUE_DIM   = RGBColor(0xDB, 0xEA, 0xFE)   # light blue tint

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Alias for clarity in slide functions
BG = BG_DARK


def box(slide, x, y, w, h, fill_color=SURFACE, line_color=BORDER, line_width=Pt(0.75)):
    """Add a rounded-rect card."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, x, y, w, h,
          size=11, bold=False, color=TEXT,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def pill(slide, text, x, y, w, h, bg_color, txt_color, size=8):
    """Colored status pill."""
    b = box(slide, x, y, w, h, fill_color=bg_color, line_color=None)
    label(slide, text, x + 0.04, y + 0.01, w - 0.08, h - 0.02,
          size=size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)
    return b


def hline(slide, x, y, w, color=BORDER, width=Pt(0.5)):
    """Thin horizontal rule."""
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = width


def status_circle(slide, letter, x, y, size=0.38, status="green"):
    color_map = {"green": GREEN_DIM, "yellow": YELLOW_DIM, "red": RED_DIM}
    txt_map   = {"green": GREEN,     "yellow": YELLOW,     "red": RED}
    s = size
    b = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(s), Inches(s))  # oval
    b.fill.solid()
    b.fill.fore_color.rgb = color_map[status]
    b.line.fill.background()
    label(slide, letter, x, y, s, s - 0.04,
          size=10, bold=True, color=txt_map[status], align=PP_ALIGN.CENTER)


def progress_bar(slide, x, y, w, h, pct, fill_color=ACCENT):
    """Background track + filled portion."""
    box(slide, x, y, w, h, fill_color=BORDER, line_color=None)
    if pct > 0:
        box(slide, x, y, w * pct / 100, h, fill_color=fill_color, line_color=None)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title / Cover
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    bg(sl)

    # Gradient accent bar left edge
    accent_bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()

    # Decorative top strip
    top_strip = sl.shapes.add_shape(1, Inches(0.18), Inches(0), SLIDE_W - Inches(0.18), Inches(0.06))
    top_strip.fill.solid()
    top_strip.fill.fore_color.rgb = SURFACE
    top_strip.line.fill.background()

    # Main title
    label(sl, "Radiological Health Program",
          0.5, 1.8, 12, 0.7,
          size=32, bold=True, color=TEXT)
    label(sl, "Modernization Initiative",
          0.5, 2.45, 12, 0.6,
          size=26, bold=False, color=ACCENT)

    # Sub-line
    label(sl, "Executive Status Report  ·  FY26  ·  June 12, 2026",
          0.5, 3.2, 10, 0.4,
          size=13, color=MUTED)

    # Overall status badge
    pill(sl, "⬤  Overall Status: YELLOW", 0.5, 3.85, 2.8, 0.38, YELLOW_DIM, YELLOW, size=10)

    hline(sl, 0.5, 4.45, 12.5, color=BORDER)

    label(sl, "Maryland Department of the Environment (MDE)  ·  PCA: 79103  ·  Fund: Radiation Control Fund",
          0.5, 4.6, 12.5, 0.35, size=9, color=MUTED)

    label(sl, "Prepared for Executive Review  ·  July 10, 2026",
          0.5, 4.95, 12.5, 0.35, size=9, color=MUTED)

    # Icon watermark area
    label(sl, "☢", 9.5, 1.5, 3.5, 3.5, size=120, color=RGBColor(0xC0, 0xD4, 0xF0), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Project Overview
# ═══════════════════════════════════════════════════════════════════════════════

def slide_overview(prs):
    sl = blank_slide(prs)
    bg(sl)

    # Header bar
    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "Project Overview", 0.3, 0.1, 8, 0.5, size=16, bold=True, color=TEXT)
    pill(sl, "Overall: YELLOW", 11.2, 0.17, 1.9, 0.36, YELLOW_DIM, YELLOW, size=9)

    # Description card
    box(sl, 0.3, 0.9, 12.7, 1.35)
    label(sl, "INITIATIVE", 0.55, 0.95, 5, 0.25, size=7, bold=True, color=MUTED)
    label(sl,
          "The Radiological Health Program (RHP) is launching a strategic modernization initiative to transition "
          "its permitting and licensing application pipeline from legacy, manual workflows to a centralized "
          "online portal and electronic payment module.",
          0.55, 1.18, 12.2, 0.9, size=10.5, color=TEXT, wrap=True)

    # Health status indicators
    label(sl, "HEALTH STATUS INDICATORS", 0.3, 2.42, 5, 0.25, size=7, bold=True, color=MUTED)

    dims = [
        ("Scope",    "G", "green"),
        ("Schedule", "G", "green"),
        ("Cost",     "G", "green"),
        ("Risk",     "G", "green"),
        ("Quality",  "R", "red"),
        ("Resources","Y", "yellow"),
    ]
    card_w = 2.05
    for i, (name, letter, status) in enumerate(dims):
        cx = 0.3 + i * (card_w + 0.07)
        box(sl, cx, 2.72, card_w, 1.05, fill_color=SURFACE2)
        status_circle(sl, letter, cx + 0.82, 2.80, status=status)
        label(sl, name, cx, 3.26, card_w, 0.3, size=9, bold=True,
              color=TEXT, align=PP_ALIGN.CENTER)

    # Yellow status rationale
    box(sl, 0.3, 3.92, 12.7, 0.92, fill_color=YELLOW_DIM, line_color=YELLOW)
    label(sl, "⚠  WHY YELLOW?", 0.55, 3.98, 5, 0.28, size=8, bold=True, color=YELLOW)
    label(sl,
          "Environment and resource changes have been identified during the reporting period. "
          "These changes are not expected to impact overall project delivery. "
          "The team is actively hiring a Tech Lead to prevent future environmental and build blockers, "
          "as key agency technical resources are engaged in other internal initiatives.",
          0.55, 4.24, 12.1, 0.55, size=9.5, color=TEXT, wrap=True)

    hline(sl, 0.3, 5.05, 12.7)
    label(sl, "Reporting Period: FY26  ·  Jun 12, 2026", 0.3, 5.12, 12.7, 0.3,
          size=8.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Sprint Schedule
# ═══════════════════════════════════════════════════════════════════════════════

def slide_sprints(prs):
    sl = blank_slide(prs)
    bg(sl)

    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "Sprint Schedule & Delivery Milestones", 0.3, 0.1, 10, 0.5, size=16, bold=True, color=TEXT)
    pill(sl, "Sprint 2: IN PROGRESS", 10.6, 0.17, 2.45, 0.36, BLUE_DIM, ACCENT, size=9)

    sprints = [
        ("Sprint 1", "Jun 17 – Jun 30, 2026", "Complete",    100, GREEN,  "Group 3–9: PDF Portal Submissions & Electronic Payments enabled"),
        ("Sprint 2", "Jul 1 – Jul 14, 2026",  "In Progress",  64, ACCENT, "Group 3–9: Enable Notifications for customers & RHP Team Inbox"),
        ("Sprint 3", "Jul 15 – Jul 28, 2026", "Not Started",   0, MUTED,  "Group 3–9: Provide option to download RHP permits"),
        ("Sprint 4", "Jul 29 – Aug 11, 2026", "Not Started",   0, MUTED,  "Group 3–9: Reject & update status (CRM) for RHP permits"),
        ("Sprint 5", "TBD",                   "Not Started",   0, MUTED,  "Group 1–2: Configuration and setup"),
        ("Sprint 6", "TBD",                   "Not Started",   0, MUTED,  "Group 3–9: Web Forms Development"),
    ]

    status_colors = {
        "Complete":    (GREEN_DIM,  GREEN),
        "In Progress": (BLUE_DIM,   ACCENT),
        "Not Started": (SURFACE2,   MUTED),
    }

    row_h   = 0.82
    row_gap = 0.04
    start_y = 0.85

    for i, (name, dates, status, pct, bar_color, goal) in enumerate(sprints):
        y = start_y + i * (row_h + row_gap)
        box(sl, 0.3, y, 12.7, row_h, fill_color=SURFACE)

        # Sprint label
        label(sl, name, 0.5, y + 0.07, 1.1, 0.3, size=11, bold=True, color=TEXT)

        # Dates
        label(sl, dates, 0.5, y + 0.42, 1.8, 0.28, size=8, color=MUTED)

        # Goal text
        label(sl, goal, 1.75, y + 0.12, 7.3, 0.55, size=9, color=TEXT, wrap=True)

        # Progress bar
        progress_bar(sl, 9.4, y + 0.22, 2.2, 0.16, pct, fill_color=bar_color)
        label(sl, f"{pct}%", 11.68, y + 0.17, 0.45, 0.26, size=8, bold=True, color=TEXT)

        # Status pill
        bg_c, txt_c = status_colors[status]
        pill(sl, status, 9.4, y + 0.48, 1.35, 0.24, bg_c, txt_c, size=7.5)

    hline(sl, 0.3, 7.15, 12.7)
    label(sl, "Near = within 2 sprints  ·  Next = planned future sprints", 0.3, 7.2, 12.7, 0.25,
          size=8, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Financial Summary
# ═══════════════════════════════════════════════════════════════════════════════

def slide_financials(prs):
    sl = blank_slide(prs)
    bg(sl)

    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "Financial Summary", 0.3, 0.1, 10, 0.5, size=16, bold=True, color=TEXT)
    pill(sl, "Cost Status: GREEN", 11.0, 0.17, 2.0, 0.36, GREEN_DIM, GREEN, size=9)

    # ── KPI cards ──────────────────────────────────────────────────────────────
    kpis = [
        ("Total Approved",  "$6.0M",   "FY25–FY26 Combined",  ACCENT, BLUE_DIM),
        ("Total Forecast",  "$3.70M",  "FY25–FY28 All Years", YELLOW, YELLOW_DIM),
        ("Total Actual",    "$1.45M",  "FY25 + FY26 Spent",   GREEN,  GREEN_DIM),
        ("Remaining",       "$1.55M",  "Balance / Variance",  PURPLE, RGBColor(0xED, 0xE9, 0xFE)),
    ]
    kw = 3.08
    for i, (lbl, val, sub, col, bg_col) in enumerate(kpis):
        cx = 0.3 + i * (kw + 0.07)
        box(sl, cx, 0.85, kw, 1.0, fill_color=bg_col, line_color=col)
        label(sl, lbl, cx + 0.12, 0.9,  kw - 0.24, 0.28, size=7.5, bold=True, color=MUTED)
        label(sl, val,  cx + 0.12, 1.12, kw - 0.24, 0.45, size=20, bold=True, color=col)
        label(sl, sub,  cx + 0.12, 1.56, kw - 0.24, 0.22, size=7.5, color=MUTED)

    # ── Table ──────────────────────────────────────────────────────────────────
    table_data = [
        ["Fiscal Year", "Approved",      "Forecast",     "Actual",       "Balance",     "Fund Source"],
        ["FY-25",       "$1,000,000",    "$500,000",     "$275,086",     "$724,914",    "Radiation Control (SF)"],
        ["FY-26",       "$2,000,000",    "$1,200,000",   "$1,177,504",   "$822,496",    "Radiation Control (SF)"],
        ["FY-27",       "$3,000,000",    "$1,468,511",   "—",            "$0",          "Special"],
        ["FY-28",       "—",             "$535,591",     "—",            "$0",          "Special"],
        ["TOTAL",       "$6,000,000",    "$3,704,102",   "$1,452,590",   "$1,547,410",  "—"],
    ]

    cols = [1.35, 1.55, 1.55, 1.55, 1.45, 2.6]
    col_x = [0.3]
    for c in cols[:-1]:
        col_x.append(col_x[-1] + c + 0.04)

    row_h   = 0.44
    start_y = 2.08

    for r, row in enumerate(table_data):
        y = start_y + r * row_h
        is_header = r == 0
        is_total  = r == len(table_data) - 1

        if is_header:
            row_bg = BLUE_DIM
        elif is_total:
            row_bg = RGBColor(0xDB, 0xEA, 0xFE)
        elif r % 2 == 0:
            row_bg = SURFACE2
        else:
            row_bg = SURFACE

        for c_idx, (cell_text, cx, cw) in enumerate(zip(row, col_x, cols)):
            box(sl, cx, y, cw, row_h, fill_color=row_bg,
                line_color=BORDER if not is_header else None, line_width=Pt(0.4))

            if is_header:
                txt_color = MUTED
                fsize, fbold = 7.5, True
            elif is_total:
                txt_color = TEXT
                fsize, fbold = 9.5, True
            elif c_idx == 4 and cell_text not in ("$0", "—"):
                txt_color = GREEN
                fsize, fbold = 9, True
            else:
                txt_color = TEXT
                fsize, fbold = 9, c_idx == 0

            align = PP_ALIGN.CENTER if c_idx == 0 else (
                PP_ALIGN.RIGHT if c_idx in (1, 2, 3, 4) else PP_ALIGN.LEFT
            )
            label(sl, cell_text, cx + 0.06, y + 0.1, cw - 0.12, row_h - 0.08,
                  size=fsize, bold=fbold, color=txt_color, align=align)

    # ── Budget utilization bars ────────────────────────────────────────────────
    label(sl, "BUDGET UTILIZATION (ACTUAL vs APPROVED)", 0.3, 4.82, 6, 0.25,
          size=7, bold=True, color=MUTED)

    util_rows = [
        ("FY-25", 27.5, GREEN),
        ("FY-26", 58.9, ACCENT),
        ("FY-27",  0.0, MUTED),
    ]
    for i, (fy, pct, col) in enumerate(util_rows):
        y = 5.1 + i * 0.45
        label(sl, fy, 0.3, y, 0.65, 0.32, size=9, color=MUTED)
        progress_bar(sl, 1.1, y + 0.06, 9.0, 0.2, pct, fill_color=col)
        label(sl, f"{pct}%", 10.2, y, 0.7, 0.32, size=9, bold=True, color=TEXT)

    hline(sl, 0.3, 7.15, 12.7)
    label(sl, "PCA: 79103  ·  Object Code: Radiation Control Fund", 0.3, 7.2, 12.7, 0.25,
          size=8, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – MVP Capabilities
# ═══════════════════════════════════════════════════════════════════════════════

def slide_mvp(prs):
    sl = blank_slide(prs)
    bg(sl)

    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "MVP Capabilities — Groups 3 through 9", 0.3, 0.1, 10, 0.5, size=16, bold=True, color=TEXT)

    capabilities = [
        ("📄", "Portal Submissions",
         "PDF application submissions through the ESC Portal for Groups 3–9.",
         ACCENT, BLUE_DIM),
        ("💳", "Electronic Payments",
         "Credit/debit card and E-check capabilities to immediately resolve clearing delays.",
         GREEN, GREEN_DIM),
        ("🔔", "Notifications",
         "One-time email confirmations for end-users upon application and payment submission.",
         YELLOW, YELLOW_DIM),
        ("📥", "RHP Team Inbox",
         "Automated notifications delivered to the MDE Business User (RHP Team) inbox.",
         PURPLE, RGBColor(0xED, 0xE9, 0xFE)),
        ("📁", "Document Management",
         "Secure view, download of applications and supporting documents — individual files or ZIP.",
         RGBColor(0xC2, 0x41, 0x0C), RGBColor(0xFF, 0xED, 0xD5)),
        ("🔐", "IAM Authentication",
         "MDE Business User authentication via Identity and Access Management (IAM) protocols.",
         RGBColor(0x08, 0x91, 0xB2), RGBColor(0xCC, 0xF0, 0xF8)),
    ]

    cols = 3
    card_w = 4.12
    card_h = 1.78
    gap_x  = 0.18
    gap_y  = 0.2
    start_x = 0.3
    start_y = 0.88

    for i, (icon, name, desc, col, bg_col) in enumerate(capabilities):
        r = i // cols
        c = i % cols
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)

        box(sl, x, y, card_w, card_h, fill_color=bg_col, line_color=col, line_width=Pt(0.6))

        # Icon circle
        circ = sl.shapes.add_shape(9, Inches(x + 0.18), Inches(y + 0.2),
                                   Inches(0.52), Inches(0.52))
        circ.fill.solid(); circ.fill.fore_color.rgb = SURFACE2
        circ.line.fill.background()
        label(sl, icon, x + 0.18, y + 0.17, 0.52, 0.52,
              size=16, align=PP_ALIGN.CENTER, color=TEXT)

        label(sl, name, x + 0.85, y + 0.22, card_w - 0.98, 0.35,
              size=11, bold=True, color=TEXT)
        label(sl, desc, x + 0.18, y + 0.82, card_w - 0.36, 0.85,
              size=9, color=MUTED, wrap=True)

    hline(sl, 0.3, 7.15, 12.7)
    label(sl, "Groups 1–2 configuration and setup planned in Sprint 5", 0.3, 7.2, 12.7, 0.25,
          size=8, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Risks & Impediments
# ═══════════════════════════════════════════════════════════════════════════════

def slide_risks(prs):
    sl = blank_slide(prs)
    bg(sl)

    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "Risks, Impediments & Action Items", 0.3, 0.1, 10, 0.5, size=16, bold=True, color=TEXT)

    # ── Risks ──────────────────────────────────────────────────────────────────
    label(sl, "ACTIVE RISKS & IMPEDIMENTS", 0.3, 0.88, 6, 0.26, size=7.5, bold=True, color=MUTED)

    risks = [
        ("HIGH",   RED,    RED_DIM,    "Environment & Build Issues",
         "Encountered environment and build issues impacting iterative delivery. "
         "Development and validation processes are established and being monitored. "
         "Mitigation: Hiring a Tech Lead to address future environmental/build blockers.",
         "Mitigate"),
        ("MEDIUM", YELLOW, YELLOW_DIM, "Workflow Finalization",
         "Non-ETS workflow approach requires finalization: future status tracking, "
         "Third-Party access (Inspectors), Liability Checks, and invoice payment features. "
         "DoIT & MDE teams reviewing enterprise solutions.",
         "Transfer"),
        ("MEDIUM", YELLOW, YELLOW_DIM, "Registration Number Automation",
         "Program must decide on folder naming convention and whether the system "
         "auto-generates sequence numbers as part of modernization, or continues the current manual process. "
         "Awaiting program decision.",
         "Mitigate"),
    ]

    for i, (level, col, bg_col, title, desc, strategy) in enumerate(risks):
        y = 1.22 + i * 1.28
        box(sl, 0.3, y, 9.5, 1.12, fill_color=bg_col, line_color=col, line_width=Pt(0.75))
        pill(sl, level, 0.45, y + 0.12, 0.72, 0.28, col, WHITE, size=7.5)
        label(sl, title, 1.3, y + 0.1, 8.0, 0.3, size=10.5, bold=True, color=TEXT)
        label(sl, desc, 0.45, y + 0.48, 9.1, 0.58, size=8.5, color=MUTED, wrap=True)

        # Strategy tag
        box(sl, 9.95, y, 3.05, 1.12, fill_color=SURFACE)
        label(sl, "RESPONSE STRATEGY", 10.1, y + 0.1, 2.7, 0.24, size=7, bold=True, color=MUTED)
        label(sl, strategy, 10.1, y + 0.38, 2.7, 0.4, size=13, bold=True, color=col, align=PP_ALIGN.LEFT)

    # ── Action Items ───────────────────────────────────────────────────────────
    label(sl, "ACTION ITEMS", 0.3, 5.12, 4, 0.26, size=7.5, bold=True, color=MUTED)

    actions = [
        (ACCENT,  "MVP Time Savings",
         "Collaborate with program team to generate high-level estimates of days saved per MVP deliverable (80/20 rule)."),
        (YELLOW,  "Liability Check",
         "Remain closely coordinated with Emil, Paul, and the DoIT team while awaiting update from State Data Officer."),
        (GREEN,   "User Testing",
         "Align with Paul to ensure user testing is seamlessly integrated into development cycles."),
    ]
    aw = 4.07
    for i, (col, title, desc) in enumerate(actions):
        x = 0.3 + i * (aw + 0.13)
        box(sl, x, 5.44, aw, 1.72, fill_color=SURFACE)
        # Top color stripe
        stripe = sl.shapes.add_shape(1, Inches(x), Inches(5.44), Inches(aw), Inches(0.055))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col; stripe.line.fill.background()
        label(sl, title, x + 0.14, 5.52, aw - 0.28, 0.3, size=10, bold=True, color=TEXT)
        label(sl, desc, x + 0.14, 5.86, aw - 0.28, 1.0, size=8.5, color=MUTED, wrap=True)

    hline(sl, 0.3, 7.15, 12.7)
    label(sl, "Risk Strategy: Mitigate = internal action  ·  Transfer = escalate / partner", 0.3, 7.2, 12.7, 0.25,
          size=8, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Successes & Key Progress
# ═══════════════════════════════════════════════════════════════════════════════

def slide_successes(prs):
    sl = blank_slide(prs)
    bg(sl)

    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "Successes & Key Progress — Reporting Period", 0.3, 0.1, 10, 0.5, size=16, bold=True, color=TEXT)

    successes = [
        ("✅", GREEN,  "MVP Release Development",
         "The team is actively working on the MVP release supporting New Applications for Groups 3 through 9. "
         "Re-certification and Renewal workflows are under review pending Liability Check clarification "
         "(Certificate of Good Standing — new system does not require Tax ID or SSN)."),
        ("⚙️", ACCENT, "ESC Page Validations & Verification",
         "The RHP Team is developing the RHP Configuration Template to govern ESC Portal page validations, "
         "ensuring data integrity and compliance across all application submissions."),
        ("🔐", PURPLE, "Authentication & Database Restructuring",
         "Work has commenced on MDE Business User authentication using Identity and Access Management (IAM) "
         "protocols, alongside the necessary restructuring of legacy RHP applications and database tables "
         "to support modern workflows."),
    ]

    for i, (icon, col, title, desc) in enumerate(successes):
        y = 0.88 + i * 1.82
        box(sl, 0.3, y, 12.7, 1.65, fill_color=SURFACE)

        # Left color bar
        bar = sl.shapes.add_shape(1, Inches(0.3), Inches(y), Inches(0.07), Inches(1.65))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()

        # Icon
        label(sl, icon, 0.55, y + 0.22, 0.65, 0.65, size=26, align=PP_ALIGN.CENTER, color=TEXT)

        label(sl, title, 1.38, y + 0.18, 11.2, 0.38, size=13, bold=True, color=TEXT)
        label(sl, desc,  1.38, y + 0.6,  11.2, 0.95, size=9.5, color=MUTED, wrap=True)

    hline(sl, 0.3, 7.15, 12.7)
    label(sl, "Sprint 1 completed successfully  ·  Sprint 2 currently in progress (Jul 1–14, 2026)",
          0.3, 7.2, 12.7, 0.25, size=8, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Executive Summary / Next Steps
# ═══════════════════════════════════════════════════════════════════════════════

def slide_summary(prs):
    sl = blank_slide(prs)
    bg(sl)

    hdr = sl.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.7))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = SURFACE; hdr.line.fill.background()
    label(sl, "Executive Summary & Next Steps", 0.3, 0.1, 10, 0.5, size=16, bold=True, color=TEXT)

    # At-a-glance status row
    statuses = [
        ("Overall",   "YELLOW",   YELLOW, YELLOW_DIM),
        ("Schedule",  "ON TRACK", GREEN,  GREEN_DIM),
        ("Cost",      "ON TRACK", GREEN,  GREEN_DIM),
        ("Delivery",  "MVP Q3",   ACCENT, BLUE_DIM),
    ]
    sw = 3.08
    for i, (lbl, val, col, bg_col) in enumerate(statuses):
        x = 0.3 + i * (sw + 0.07)
        box(sl, x, 0.85, sw, 0.82, fill_color=bg_col, line_color=col)
        label(sl, lbl, x + 0.12, 0.9, sw - 0.24, 0.26, size=7.5, bold=True, color=MUTED)
        label(sl, val, x + 0.12, 1.1, sw - 0.24, 0.42, size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Key takeaways
    label(sl, "KEY TAKEAWAYS", 0.3, 1.86, 6, 0.26, size=7.5, bold=True, color=MUTED)

    takeaways = [
        (GREEN,  "On Budget",
         "Actual spend of $1.45M is tracking below forecast. Remaining balance of $1.55M provides healthy fiscal buffer."),
        (GREEN,  "Delivery On Track",
         "Sprint 1 completed. Sprint 2 in progress. Core MVP capabilities on schedule for Groups 3–9."),
        (RED,    "Quality Attention Needed",
         "Build/environment issues were encountered in Sprint 1. Tech Lead hiring in progress to prevent recurrence."),
        (YELLOW, "Decisions Pending",
         "Workflow finalization and registration number automation require program decisions to maintain velocity."),
    ]

    tw = 3.0
    th = 2.3
    for i, (col, title, desc) in enumerate(takeaways):
        x = 0.3 + i * (tw + 0.12)
        box(sl, x, 2.2, tw, th, fill_color=SURFACE)
        stripe = sl.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(tw), Inches(0.06))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col; stripe.line.fill.background()
        label(sl, title, x + 0.14, 2.32, tw - 0.28, 0.32, size=10.5, bold=True, color=TEXT)
        label(sl, desc,  x + 0.14, 2.7,  tw - 0.28, 1.7, size=9, color=MUTED, wrap=True)

    # Next steps
    label(sl, "IMMEDIATE NEXT STEPS", 0.3, 4.72, 6, 0.26, size=7.5, bold=True, color=MUTED)

    next_steps = [
        (ACCENT, "Complete Sprint 2",       "Jul 1–14, 2026",    "Enable customer notifications & RHP Team inbox"),
        (ACCENT, "Finalize Workflow",       "Ongoing",           "Resolve non-ETS workflow with DoIT & MDE teams"),
        (YELLOW, "State Data Officer Input","Awaiting",          "Liability Check decision needed from Emil/Paul/DoIT"),
        (GREEN,  "Hire Tech Lead",          "Active Recruiting", "Address build/environment blockers proactively"),
        (GREEN,  "User Testing Alignment",  "Sprint 2–3",        "Coordinate with Paul to integrate testing into dev cycle"),
    ]

    for i, (col, step, timeline, detail) in enumerate(next_steps):
        y = 5.04 + i * 0.42
        box(sl, 0.3, y, 12.7, 0.38, fill_color=SURFACE2 if i % 2 == 0 else SURFACE)
        circ2 = sl.shapes.add_shape(9, Inches(0.42), Inches(y + 0.1), Inches(0.18), Inches(0.18))
        circ2.fill.solid(); circ2.fill.fore_color.rgb = col; circ2.line.fill.background()
        label(sl, step,     0.72, y + 0.06, 3.2,  0.28, size=9.5, bold=True, color=TEXT)
        label(sl, timeline, 3.95, y + 0.06, 1.9,  0.28, size=9,   color=col)
        label(sl, detail,   5.92, y + 0.06, 7.0,  0.28, size=9,   color=MUTED)

    hline(sl, 0.3, 7.15, 12.7)
    label(sl, "RHP Modernization Initiative  ·  MDE  ·  Prepared for Executive Review  ·  July 10, 2026",
          0.3, 7.2, 12.7, 0.25, size=8, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_title(prs)
    slide_overview(prs)
    slide_sprints(prs)
    slide_financials(prs)
    slide_mvp(prs)
    slide_risks(prs)
    slide_successes(prs)
    slide_summary(prs)

    out = "dashboard/RHP_Executive_Status.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({prs.slides.__len__()} slides)")

if __name__ == "__main__":
    build()
